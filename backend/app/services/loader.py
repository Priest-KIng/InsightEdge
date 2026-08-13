from __future__ import annotations

import hashlib
from io import BytesIO
from pathlib import Path
import re
import statistics
from typing import Iterable

from bs4 import BeautifulSoup
from docx import Document
from ebooklib import ITEM_DOCUMENT, epub
import pandas as pd
from PIL import Image
import pdfplumber
from pptx import Presentation
from pypdf import PdfReader
import pytesseract
import structlog

from app.services.document_model import DocumentSegment

try:
    from pdf2image import convert_from_path
except Exception:
    convert_from_path = None

logger = structlog.get_logger(__name__)

SUPPORTED_EXTENSIONS = {
    ".txt",
    ".md",
    ".pdf",
    ".docx",
    ".csv",
    ".xlsx",
    ".html",
    ".htm",
    ".epub",
    ".pptx",
}


class DocumentLoadError(Exception):
    pass


def iter_supported_files(path: Path) -> Iterable[Path]:
    if path.is_file() and path.suffix.lower() in SUPPORTED_EXTENSIONS:
        yield path
        return

    if path.is_dir():
        for file in path.rglob("*"):
            if file.is_file() and file.suffix.lower() in SUPPORTED_EXTENSIONS:
                yield file


def load_text(file_path: Path) -> str:
    suffix = file_path.suffix.lower()

    try:
        if suffix in {".txt", ".md"}:
            return file_path.read_text(encoding="utf-8", errors="ignore")
        if suffix == ".pdf":
            return _load_pdf_text(file_path)
        if suffix == ".docx":
            return _load_docx_text(file_path)
        if suffix in {".csv", ".xlsx"}:
            return _load_tabular_text(file_path)
        if suffix in {".html", ".htm"}:
            return _load_html_text(file_path)
        if suffix == ".epub":
            return _load_epub_text(file_path)
        if suffix == ".pptx":
            return _load_pptx_text(file_path)
    except Exception as exc:
        raise DocumentLoadError(f"Failed to parse {file_path}: {exc}") from exc

    raise DocumentLoadError(f"Unsupported file type: {file_path.suffix}")


def load_structured(file_path: Path) -> tuple[str, list[DocumentSegment]]:
    """Load text plus structural blocks without breaking the legacy text API."""
    text = load_text(file_path)
    return text, _segment_text(text, file_path.suffix.lower())


def _segment_text(text: str, suffix: str) -> list[DocumentSegment]:
    segments: list[DocumentSegment] = []
    current_lines: list[str] = []
    current_start: int | None = None
    cursor = 0
    page_number: int | None = None
    slide_number: int | None = None
    section_title: str | None = None
    block_type = "table" if suffix in {".csv", ".xlsx"} else "paragraph"
    ocr_used = False
    table_used = suffix in {".csv", ".xlsx"}

    marker_pattern = re.compile(
        r"^\[(?:PAGE\s+(?P<page>\d+)(?:\s+(?P<page_kind>TEXT|OCR|TABLE.*|IMAGE.*))?|SLIDE\s+(?P<slide>\d+))\]",
        flags=re.IGNORECASE,
    )

    def flush(end_position: int) -> None:
        nonlocal current_lines, current_start, block_type, ocr_used, table_used
        value = "\n".join(current_lines).strip()
        if value:
            segments.append(
                DocumentSegment(
                    text=value,
                    block_type=block_type,
                    section_title=section_title,
                    page_number=page_number,
                    slide_number=slide_number,
                    start_char=current_start,
                    end_char=end_position,
                    ocr_used=ocr_used or "[OCR]" in value.upper(),
                    table_used=table_used or "|" in value or "TABLE" in value.upper(),
                ),
            )
        current_lines = []
        current_start = None
        block_type = "table" if suffix in {".csv", ".xlsx"} else "paragraph"
        ocr_used = False
        table_used = suffix in {".csv", ".xlsx"}

    for raw_line in text.splitlines(keepends=True):
        line = raw_line.rstrip("\r\n")
        stripped = line.strip()
        marker = marker_pattern.match(stripped)
        if marker:
            flush(cursor)
            if marker.group("page"):
                page_number = int(marker.group("page"))
                page_kind = (marker.group("page_kind") or "").upper()
                ocr_used = "OCR" in page_kind
                table_used = "TABLE" in page_kind
                block_type = "table" if table_used else ("ocr" if ocr_used else "page")
            if marker.group("slide"):
                slide_number = int(marker.group("slide"))
                block_type = "slide"
            cursor += len(raw_line)
            continue
        if not stripped:
            flush(cursor)
            cursor += len(raw_line)
            continue

        heading = re.match(r"^#{1,6}\s+(.+)$", stripped)
        if heading:
            flush(cursor)
            section_title = heading.group(1).strip()
            heading_start = cursor
            segments.append(
                DocumentSegment(
                    text=section_title,
                    block_type="heading",
                    section_title=section_title,
                    page_number=page_number,
                    slide_number=slide_number,
                    start_char=heading_start,
                    end_char=heading_start + len(line),
                    ocr_used=ocr_used,
                    table_used=table_used,
                ),
            )
            current_start = None
            current_lines = []
            block_type = "paragraph"
        else:
            if current_start is None:
                current_start = cursor
            if "|" in stripped or stripped.upper().startswith("[TABLE"):
                block_type = "table"
                table_used = True
            elif stripped.startswith(("- ", "* ", "• ")):
                block_type = "bullet"
            elif block_type == "heading":
                block_type = "paragraph"
            current_lines.append(stripped)
            if "[OCR]" in stripped.upper():
                ocr_used = True
        cursor += len(raw_line)

    flush(len(text))
    if segments:
        return segments
    if text.strip():
        return [DocumentSegment(text=text.strip(), start_char=0, end_char=len(text))]
    return []


def _load_tabular_text(file_path: Path) -> str:
    suffix = file_path.suffix.lower()
    if suffix == ".csv":
        frame = pd.read_csv(file_path)
        return "[TABLE CSV]\n" + frame.to_csv(index=False)

    sheets = pd.read_excel(file_path, sheet_name=None)
    parts: list[str] = []
    for sheet_name, frame in sheets.items():
        parts.append(f"[TABLE SHEET {sheet_name}]")
        parts.append(frame.to_csv(index=False))
    return "\n\n".join(parts).strip()


def _load_html_text(file_path: Path) -> str:
    html = file_path.read_text(encoding="utf-8", errors="ignore")
    soup = BeautifulSoup(html, "html.parser")
    parts: list[str] = []
    for element in soup.find_all(["h1", "h2", "h3", "h4", "h5", "h6", "table", "p", "li"]):
        value = element.get_text(" ", strip=True)
        if not value:
            continue
        if element.name and element.name.startswith("h"):
            parts.append(f"{'#' * int(element.name[1])} {value}")
        elif element.name == "table":
            rows = []
            for row in element.find_all("tr"):
                cells = [cell.get_text(" ", strip=True) for cell in row.find_all(["th", "td"])]
                if cells:
                    rows.append(" | ".join(cells))
            if rows:
                parts.append("[TABLE HTML]\n" + "\n".join(rows))
        else:
            parts.append(value)
    return "\n\n".join(parts).strip()


def _load_epub_text(file_path: Path) -> str:
    book = epub.read_epub(str(file_path))
    parts: list[str] = []
    for item in book.get_items():
        if item.get_type() != ITEM_DOCUMENT:
            continue
        soup = BeautifulSoup(item.get_content(), "html.parser")
        text = soup.get_text(separator="\n", strip=True)
        if text:
            title = item.get_name()
            parts.append(f"[EPUB {title}]\n{text}")
    return "\n\n".join(parts).strip()


def _load_pptx_text(file_path: Path) -> str:
    deck = Presentation(str(file_path))
    slide_sections: list[str] = []
    for slide_index, slide in enumerate(deck.slides):
        lines: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                cleaned = str(shape.text).strip()
                if cleaned:
                    lines.append(cleaned)
            if getattr(shape, "has_table", False):
                table_lines: list[str] = []
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    if any(cells):
                        table_lines.append(" | ".join(cells))
                if table_lines:
                    lines.append("\n".join(table_lines))
        if lines:
            slide_sections.append(f"[SLIDE {slide_index + 1}]\n" + "\n".join(lines))
    return "\n\n".join(slide_sections).strip()


def _load_pdf_text(file_path: Path) -> str:
    text_sections: list[str] = []
    try:
        with pdfplumber.open(str(file_path)) as pdf:
            for page_index, page in enumerate(pdf.pages):
                page_text = (page.extract_text() or "").strip()
                if page_text:
                    page_text = _mark_pdf_headings(page, page_text)
                    text_sections.append(f"[PAGE {page_index + 1} TEXT]\n{page_text}")
                else:
                    ocr_page_text = _ocr_full_pdf_page(file_path, page_index + 1)
                    if ocr_page_text:
                        text_sections.append(f"[PAGE {page_index + 1} OCR]\n{ocr_page_text}")
                    else:
                        text_sections.append(f"[PAGE {page_index + 1}] OCR fallback unavailable or no readable text.")

                tables = page.extract_tables() or []
                for table_index, table in enumerate(tables):
                    normalized_rows: list[str] = []
                    for row in table:
                        cells = [str(cell).strip() if cell is not None else "" for cell in row]
                        if any(cells):
                            normalized_rows.append(" | ".join(cells))

                    if normalized_rows:
                        table_blob = "\n".join(normalized_rows)
                        text_sections.append(f"[PAGE {page_index + 1} TABLE {table_index + 1}]\n{table_blob}")
    except Exception as exc:
        logger.warning("pdfplumber_extraction_failed", file=str(file_path), error=str(exc))

    try:
        reader = PdfReader(str(file_path))
        for page_index, page in enumerate(reader.pages):
            page_images = list(page.images)
            if page_images:
                text_sections.append(f"[PAGE {page_index + 1}] Found {len(page_images)} embedded image(s).")

            for image_index, image in enumerate(page_images):
                image_note = _extract_image_ocr(image.data, page_index + 1, image_index + 1)
                if image_note:
                    text_sections.append(image_note)
    except Exception as exc:
        logger.warning("pypdf_image_extraction_failed", file=str(file_path), error=str(exc))

    final_text = "\n\n".join([section for section in text_sections if section.strip()]).strip()
    if not final_text:
        raise DocumentLoadError(
            "No text extracted from PDF. Ensure tesseract and poppler are installed for OCR fallback.",
        )

    return final_text


def _mark_pdf_headings(page, page_text: str) -> str:
    """Add heading markers to PDF lines using layout/font evidence when available."""
    try:
        words = page.extract_words(extra_attrs=["fontname", "size"], use_text_flow=True)
    except Exception:
        return page_text
    if not words:
        return page_text

    lines: dict[tuple[int, int], list[dict[str, object]]] = {}
    for word in words:
        top = float(word.get("top", 0.0))
        bottom = float(word.get("bottom", top))
        key = (round(top), round(bottom))
        lines.setdefault(key, []).append(word)

    sizes = [float(word.get("size", 0.0)) for word in words if word.get("size") is not None]
    median_size = statistics.median(sizes) if sizes else 0.0
    candidates: set[str] = set()
    for line_words in lines.values():
        line_words.sort(key=lambda item: float(item.get("x0", 0.0)))
        value = " ".join(str(item.get("text", "")).strip() for item in line_words).strip()
        if not value or len(value) > 120 or len(line_words) > 18:
            continue
        max_size = max(float(item.get("size", 0.0)) for item in line_words)
        fonts = " ".join(str(item.get("fontname", "")).lower() for item in line_words)
        looks_academic = bool(re.match(r"^(?:[IVXLCDM]+|[A-Z])\\.\\s+[A-Z]", value))
        looks_styled = median_size > 0 and max_size >= median_size * 1.12
        looks_bold = any(token in fonts for token in ("bold", "black", "demi", "semibold"))
        if (looks_academic or looks_styled or looks_bold) and not re.match(
            r"^(?:fig(?:ure)?|table)\\b",
            value,
            flags=re.IGNORECASE,
        ):
            candidates.add(" ".join(value.split()).lower())

    if not candidates:
        return page_text

    marked_lines: list[str] = []
    for line in page_text.splitlines():
        normalized = " ".join(line.split()).lower()
        if normalized in candidates and not line.lstrip().startswith("#"):
            marked_lines.append(f"# {line.strip()}")
        else:
            marked_lines.append(line)
    return "\n".join(marked_lines)


def _load_docx_text(file_path: Path) -> str:
    doc = Document(str(file_path))
    text_sections: list[str] = []

    paragraph_lines: list[str] = []
    for paragraph in doc.paragraphs:
        value = paragraph.text.strip()
        if not value:
            continue
        style = str(getattr(paragraph.style, "name", ""))
        if "heading" in style.lower():
            level = re.search(r"(\d+)$", style)
            prefix = "#" * int(level.group(1)) if level else "#"
            paragraph_lines.append(f"{prefix} {value}")
        else:
            paragraph_lines.append(value)
    if paragraph_lines:
        text_sections.append("\n".join(paragraph_lines))

    for table_index, table in enumerate(doc.tables, start=1):
        rows: list[str] = []
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            if any(cells):
                rows.append(" | ".join(cells))
        if rows:
            text_sections.append(f"[TABLE DOCX {table_index}]\n" + "\n".join(rows))

    seen_images: set[str] = set()
    image_index = 0
    for rel in doc.part.rels.values():
        if "image" not in rel.reltype:
            continue
        image_bytes = rel.target_part.blob
        digest = hashlib.sha256(image_bytes).hexdigest()
        if digest in seen_images:
            continue
        seen_images.add(digest)
        image_index += 1
        image_note = _extract_docx_image_ocr(image_bytes, image_index)
        if image_note:
            text_sections.append(image_note)

    return "\n\n".join([section for section in text_sections if section.strip()]).strip()


def _ocr_full_pdf_page(file_path: Path, page_number: int) -> str:
    if convert_from_path is None:
        return ""

    try:
        images = convert_from_path(
            str(file_path),
            first_page=page_number,
            last_page=page_number,
            dpi=300,
        )
    except Exception as exc:
        logger.warning(
            "pdf_full_page_ocr_render_failed",
            file=str(file_path),
            page=page_number,
            error=str(exc),
        )
        return ""

    if not images:
        return ""

    try:
        image = images[0]
        if image.mode != "RGB":
            image = image.convert("RGB")
        return pytesseract.image_to_string(image).strip()
    except Exception as exc:
        logger.warning(
            "pdf_full_page_ocr_failed",
            file=str(file_path),
            page=page_number,
            error=str(exc),
        )
        return ""


def _extract_image_ocr(image_bytes: bytes, page_number: int, image_number: int) -> str | None:
    try:
        with Image.open(BytesIO(image_bytes)) as img:
            if img.mode != "RGB":
                img = img.convert("RGB")
            ocr_text = pytesseract.image_to_string(img).strip()
    except Exception as exc:
        return f"[PAGE {page_number} IMAGE {image_number}] Embedded image present, OCR unavailable or failed: {exc}"

    if not ocr_text:
        return f"[PAGE {page_number} IMAGE {image_number}] Embedded image present, no readable text detected."

    return f"[PAGE {page_number} IMAGE {image_number} OCR]\n{ocr_text}"


def _extract_docx_image_ocr(image_bytes: bytes, image_number: int) -> str | None:
    try:
        with Image.open(BytesIO(image_bytes)) as img:
            if img.mode != "RGB":
                img = img.convert("RGB")
            ocr_text = pytesseract.image_to_string(img).strip()
    except Exception as exc:
        return f"[DOCX IMAGE {image_number}] Embedded image present, OCR unavailable or failed: {exc}"

    if not ocr_text:
        return f"[DOCX IMAGE {image_number}] Embedded image present, no readable text detected."

    return f"[DOCX IMAGE {image_number} OCR]\n{ocr_text}"
